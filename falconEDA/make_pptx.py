

from utils import *

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

import io


# def create_pptx(df, filename="EDA_Presentation.pptx"):
#     # Create a new PowerPoint presentation
#     prs = Presentation()

#     # Set slide dimensions to Widescreen (13.33 inches x 7.5 inches)
#     prs.slide_width = Inches(13.33)
#     prs.slide_height = Inches(7.5)

#     slide_layout = prs.slide_layouts[5]  # layout
    
#     for col_name in df.columns:
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         title.text = f"Analysis of: {col_name}"
        
#         text_frame = title.text_frame
#         for paragraph in text_frame.paragraphs:
#             paragraph.alignment = PP_ALIGN.LEFT  # Left align

#         col_type = df[col_name].dtype
        
#         if col_type == 'object':  # Categorical column
#             plot_data = bar_chart_data(df, col_name, top_n_rows=5)
            
#             # Generate the bar chart
#             fig_width, fig_height = 6.5, 4.5
#             fig, ax = plt.subplots(figsize=(fig_width, fig_height))
#             ax.barh(plot_data[col_name], plot_data['Occurrences'], color='skyblue')
#             ax.set_xlabel("Occurrences")
#             ax.set_title(f"Top 5 Values in {col_name}")
#             plt.tight_layout()
            
#             # Save the plot to a buffer
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi = 400)
#             img_buffer.seek(0)
#             plt.close()
            
#             left = Inches(0.05)
#             top = Inches(2.5)
#             # Add the chart to the slide
#             slide.shapes.add_picture(img_buffer, left, top, width=Inches(fig_width), height=Inches(fig_height))
        
#         elif col_type in ['int64', 'float64']:  # Numeric column
#             fig_width, fig_height = 6.5, 4.5

#             fig, ax = plt.subplots(figsize=(fig_width, fig_height))
#             ax.hist(df[col_name].dropna(), bins=20, color='coral', edgecolor='black')
#             ax.set_title(f"Distribution of {col_name}")
#             ax.set_xlabel(col_name)
#             ax.set_ylabel("Frequency")
#             plt.tight_layout()
            
#             # Save the plot to a buffer
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi = 400)
#             img_buffer.seek(0)
#             plt.close()
            
#             #Add the chart to the slide
#             left = Inches(0.05)
#             top = Inches(2.5)
#             slide.shapes.add_picture(img_buffer, left, top, width=Inches(fig_width), height=Inches(fig_height))
        
#     # Save the presentation
#     prs.save(filename)
#     return filename



#df = pd.read_csv(r'D:\Documents\Python\CREATE_DATA\winemag-data_first150k.csv', engine="c", low_memory=False)
#pptx = create_pptx(df)



# def create_pptx_v2(df):
#     ppt_file = "EDA_Presentation.pptx"

#     prs = Presentation()

#     # Set slide dimensions to Widescreen (13.33 inches x 7.5 inches)
#     prs.slide_width = Inches(13.33)
#     prs.slide_height = Inches(7.5)

#     slide_layout = prs.slide_layouts[5]  # layout
    
#     for col_name in df.columns:
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         title.text = f"Analysis of: {col_name}"
        
#         text_frame = title.text_frame
#         for paragraph in text_frame.paragraphs:
#             paragraph.alignment = PP_ALIGN.LEFT  # Left align

#         col_type = df[col_name].dtype
        
#         if col_type == 'object':  # Categorical column
#             plot_data = bar_chart_data(df, col_name, top_n_rows=5)
            
#             # Generate the bar chart
#             fig_width, fig_height = 6.5, 4.5
#             fig, ax = plt.subplots(figsize=(fig_width, fig_height))
#             ax.barh(plot_data[col_name], plot_data['Occurrences'], color='skyblue')
#             ax.set_xlabel("Occurrences")
#             ax.set_title(f"Top 5 Values in {col_name}")
#             plt.tight_layout()
            
#             # Save the plot to a buffer
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi = 400)
#             img_buffer.seek(0)
#             plt.close()
            
#             left = Inches(0.05)
#             top = Inches(2.5)
#             # Add the chart to the slide
#             slide.shapes.add_picture(img_buffer, left, top, width=Inches(fig_width), height=Inches(fig_height))
        
#         elif col_type in ['int64', 'float64']:  # Numeric column
#             fig_width, fig_height = 6.5, 4.5

#             fig, ax = plt.subplots(figsize=(fig_width, fig_height))
#             ax.hist(df[col_name].dropna(), bins=20, color='coral', edgecolor='black')
#             ax.set_title(f"Distribution of {col_name}")
#             ax.set_xlabel(col_name)
#             ax.set_ylabel("Frequency")
#             plt.tight_layout()
            
#             # Save the plot to a buffer
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi = 400)
#             img_buffer.seek(0)
#             plt.close()
            
#             #Add the chart to the slide
#             left = Inches(0.05)
#             top = Inches(2.5)
#             slide.shapes.add_picture(img_buffer, left, top, width=Inches(fig_width), height=Inches(fig_height))
        
#     # Save the presentation
#     prs.save(ppt_file)

#     with open(ppt_file, "rb") as f:
#         ppt_bytes = f.read()
#     return ppt_bytes



def create_pptx_v3(df):
    ppt_file = "EDA_Presentation_v3.pptx"

    POWERPOINT_PIXEL_PER_INCH = 96    # DO NOT CHANGE
    PIXEL_PER_INCH = 350              # Resolution of our images
    
    plot_width_pixels = 800
    plot_height_pixels = 450

    width_inches = plot_width_pixels / POWERPOINT_PIXEL_PER_INCH 
    height_inches = plot_height_pixels /POWERPOINT_PIXEL_PER_INCH 

    # Create Presentation
    prs = Presentation()
    slide_width = Inches(13.33); prs.slide_width = slide_width
    slide_height = Inches(7.5) ; prs.slide_height = slide_height
    
    #center_x = slide_width / 2
    #center_y = slide_height / 2
    #left = center_x - (width_inches / 2)
    #top = center_y - (height_inches / 2)

    

    slide_layout = prs.slide_layouts[5]     # Blank layout
    BAR_COLORS = ["#826fc2","#001f80","#4d9b1e","#f865c6","#ecd378","#ba004c","#8f4400","#f65656"]*10000

    # Loop through each column in the DataFrame
    for count, col_name in enumerate(df.columns):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"{col_name}"

        # Left-align the title
        text_frame = title.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT

        col_type = df[col_name].dtype

        # Categorical column
        if col_type == 'object':
            plot_data = bar_chart_data(df, col_name, top_n_rows=5)

            # Generate Altair bar chart
            chart = alt.Chart(plot_data).mark_bar().encode(
                x=alt.X("Occurrences:Q", axis=alt.Axis(format='d', title="Occurrences")),
                y=alt.Y(col_name, sort='-x', title=None),
                color=alt.value(BAR_COLORS[count])
            ).properties(width=plot_width_pixels, height=plot_height_pixels)   # <---- Bar Graph Dimenstions

            # Save image
            img_buffer = io.BytesIO()
            chart.save(img_buffer, format="png", ppi = PIXEL_PER_INCH)
            img_buffer.seek(0)

            left = Inches(2.5)
            top = Inches(2.0)
            
            slide.shapes.add_picture(img_buffer, left, top, width = Inches(width_inches) )

        # Numeric column
        elif col_type in ['int64', 'float64']:

            plot_width_pixels = 700
            plot_height_pixels = 350

            left = Inches(2.5)
            top = Inches(1.5)

            chart = boxplot_histogram(df, col_name, bar_color=BAR_COLORS[count], WIDTH=plot_width_pixels, HEIGHT=plot_height_pixels)

            # Save image
            img_buffer = io.BytesIO()
            chart.save(img_buffer, format="png", ppi = PIXEL_PER_INCH)
            img_buffer.seek(0)

            slide.shapes.add_picture( img_buffer, left, top,  width=Inches(width_inches)  )


    prs.save(ppt_file)

    # Return PowerPoint file as bytes
    with open(ppt_file, "rb") as f:
        ppt_bytes = f.read()
    return ppt_bytes







